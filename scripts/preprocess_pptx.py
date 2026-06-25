#!/usr/bin/env python3
"""
PPTX Preprocessing Pipeline for SE Requirements Decomposition.

Extracts text, tables, speaker notes, and embedded images from PowerPoint
presentations — commonly used for architecture overviews, design reviews,
and customer requirement presentations.

Usage:
    python preprocess_pptx.py --input architecture-review.pptx --output ./preprocessed/
    python preprocess_pptx.py --input ./docs/inputs/ --output ./preprocessed/ --batch
"""

import argparse
import json
import logging
import os
import sys
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

logger = logging.getLogger("pptx_preprocess")
logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")


# ---------------------------------------------------------------------------
# Data structures
# ---------------------------------------------------------------------------

@dataclass
class SlideInfo:
    """Metadata for one slide."""
    number: int
    title: str = ""
    text_length: int = 0
    has_notes: bool = False
    has_table: bool = False
    has_image: bool = False


@dataclass
class PptxManifest:
    """Output manifest for one presentation."""
    source: str
    total_slides: int
    slides_with_notes: int
    slides_with_tables: int
    total_images: int
    slides: list[SlideInfo] = field(default_factory=list)
    issues: list[str] = field(default_factory=list)


# ---------------------------------------------------------------------------
# Core extraction
# ---------------------------------------------------------------------------

def extract_with_markitdown(filepath: str, output_dir: str) -> tuple[str, list[str]]:
    """
    Use markitdown to convert PPTX to Markdown.
    Returns (markdown_text, issues).
    """
    issues: list[str] = []

    # Try markitdown CLI first
    try:
        import subprocess
        result = subprocess.run(
            [sys.executable, "-m", "markitdown", filepath],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0 and result.stdout.strip():
            logger.info(f"markitdown: extracted {len(result.stdout)} chars")
            return result.stdout, issues
        elif result.stderr:
            issues.append(f"markitdown warning: {result.stderr[:200]}")
            if result.stdout.strip():
                return result.stdout, issues
    except FileNotFoundError:
        issues.append("markitdown not installed. Run: pip install \"markitdown[pptx]\"")
    except Exception as exc:
        issues.append(f"markitdown failed: {exc}")

    # Fallback: python-pptx manual extraction
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        md_parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            md_parts.append(f"\n## Slide {i}\n")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            md_parts.append(f"{text}\n")
                if shape.has_table:
                    table = shape.table
                    md = _pptx_table_to_markdown(table)
                    md_parts.append(f"\n{md}\n")
            # Speaker notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    md_parts.append(f"\n> **Notes:** {notes}\n")

        text = "\n".join(md_parts)
        logger.info(f"python-pptx fallback: extracted {len(text)} chars")
        return text, issues

    except ImportError:
        issues.append("python-pptx not installed. Run: pip install python-pptx")
        return "", issues
    except Exception as exc:
        issues.append(f"python-pptx fallback failed: {exc}")
        return "", issues


def _pptx_table_to_markdown(table) -> str:
    """Convert python-pptx Table to GitHub-flavored Markdown."""
    rows = []
    for row in table.rows:
        cells = [cell.text.replace("\n", " ").replace("|", "\\|").strip() for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ""

    max_cols = max(len(r) for r in rows)
    for r in rows:
        while len(r) < max_cols:
            r.append("")

    # Trim trailing empty columns
    while max_cols > 1 and all(r[-1] == "" for r in rows):
        for r in rows:
            r.pop()
        max_cols -= 1

    lines = [
        "| " + " | ".join(rows[0]) + " |",
        "| " + " | ".join(["---"] * max_cols) + " |",
    ]
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")

    return "\n".join(lines) + "\n"


def extract_slide_metadata(filepath: str) -> tuple[list[SlideInfo], int, list[str]]:
    """Extract per-slide metadata using python-pptx."""
    issues: list[str] = []
    slides: list[SlideInfo] = []
    image_count = 0

    try:
        from pptx import Presentation
        prs = Presentation(filepath)

        for i, slide in enumerate(prs.slides, 1):
            info = SlideInfo(number=i)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if not info.title and shape.is_placeholder and hasattr(shape, 'placeholder_format'):
                        if shape.placeholder_format.type == 1:  # TITLE
                            info.title = text[:80]
                    info.text_length += len(text)
                if shape.has_table:
                    info.has_table = True
                if shape.shape_type == 13:  # PICTURE
                    info.has_image = True
                    image_count += 1
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                info.has_notes = bool(notes_text)
            slides.append(info)

        return slides, image_count, issues

    except ImportError:
        return [], 0, ["python-pptx not installed"]
    except Exception as exc:
        return [], 0, [f"Metadata extraction failed: {exc}"]


def extract_images_from_pptx(filepath: str, output_dir: str) -> int:
    """Extract embedded images from .pptx ZIP."""
    import zipfile

    images_dir = os.path.join(output_dir, "images")
    os.makedirs(images_dir, exist_ok=True)
    count = 0
    image_exts = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".emf", ".wmf", ".svg"}

    try:
        with zipfile.ZipFile(filepath, "r") as z:
            for name in z.namelist():
                ext = Path(name).suffix.lower()
                if ext in image_exts and "media" in name.lower():
                    data = z.read(name)
                    if len(data) < 100:
                        continue
                    out_path = os.path.join(images_dir, Path(name).name)
                    with open(out_path, "wb") as f:
                        f.write(data)
                    count += 1
    except Exception as exc:
        logger.warning(f"Image extraction failed: {exc}")

    if count > 0:
        logger.info(f"Extracted {count} image(s) → {images_dir}")
    return count


# ---------------------------------------------------------------------------
# Quality verification
# ---------------------------------------------------------------------------

def verify_pptx_quality(manifest: PptxManifest) -> dict[str, Any]:
    """Assess PPTX extraction quality."""
    checks: list[dict[str, Any]] = []
    score = 100.0

    # Check 1: Empty slides
    empty_slides = [s for s in manifest.slides if s.text_length < 10]
    if empty_slides:
        penalty = len(empty_slides) * (100.0 / max(len(manifest.slides), 1))
        score -= min(penalty, 30)
        checks.append({
            "dimension": "empty_slides",
            "status": "warn" if len(empty_slides) > len(manifest.slides) * 0.3 else "info",
            "detail": f"{len(empty_slides)}/{len(manifest.slides)} slides have almost no text — may be image-only or failed extraction",
            "slides": [s.number for s in empty_slides],
        })
    else:
        checks.append({"dimension": "empty_slides", "status": "pass", "detail": "All slides have text content"})

    # Check 2: Speaker notes
    if manifest.slides_with_notes == 0 and manifest.total_slides > 5:
        checks.append({
            "dimension": "speaker_notes",
            "status": "info",
            "detail": "No speaker notes found — may miss rationale/decisions not on slides",
        })
    elif manifest.slides_with_notes > 0:
        checks.append({
            "dimension": "speaker_notes",
            "status": "pass",
            "detail": f"{manifest.slides_with_notes}/{manifest.total_slides} slides have speaker notes",
        })

    # Check 3: Tables
    if manifest.slides_with_tables > 0:
        checks.append({"dimension": "tables", "status": "info", "detail": f"{manifest.slides_with_tables} slide(s) contain tables"})

    # Check 4: Image-heavy slides
    image_slides = [s for s in manifest.slides if s.has_image]
    if image_slides:
        checks.append({
            "dimension": "images",
            "status": "info",
            "detail": f"{len(image_slides)} slides contain images — use Read tool to check diagrams",
        })

    if manifest.issues:
        score -= len(manifest.issues) * 15
        checks.append({"dimension": "issues", "status": "warn", "detail": f"{len(manifest.issues)} issue(s): {'; '.join(manifest.issues[:3])}"})

    if any(c["status"] == "warn" for c in checks):
        score = min(score, 85)

    return {"overall_score": round(max(score, 0), 1), "checks": checks}


def print_pptx_quality(report: dict[str, Any]) -> None:
    icon = "[PASS]" if report["overall_score"] >= 90 else ("[WARN]" if report["overall_score"] >= 70 else "[FAIL]")
    print(f"""
======================================================================
  PPTX Extraction Quality Report
======================================================================
  Overall:    {icon} {report['overall_score']:.0f}/100
----------------------------------------------------------------------""")
    for c in report["checks"]:
        ic = {"pass": "[OK]  ", "warn": "[WARN]", "info": "[INFO]"}.get(c["status"], "      ")
        print(f"  {ic} {c['detail'][:67]}")
    print("======================================================================\n")


# ---------------------------------------------------------------------------
# Main pipeline
# ---------------------------------------------------------------------------

def process_pptx(
    filepath: str,
    output_dir: str,
    *,
    extract_images: bool = True,
    verify: bool = False,
) -> PptxManifest:
    """Run full preprocessing on a single .pptx file."""
    os.makedirs(output_dir, exist_ok=True)

    # --- Step 1: Extract text via markitdown (+ python-pptx fallback) ---
    md_text, issues = extract_with_markitdown(filepath, output_dir)
    md_path = os.path.join(output_dir, "full-text.md")
    with open(md_path, "w", encoding="utf-8") as f:
        f.write(f"<!-- Source: {os.path.basename(filepath)} -->\n\n")
        f.write(md_text)
    logger.info(f"Full text → {md_path}")

    # --- Step 2: Per-slide metadata ---
    slides, img_count_meta, meta_issues = extract_slide_metadata(filepath)
    issues.extend(meta_issues)

    # --- Step 3: Extract embedded images ---
    image_count = 0
    if extract_images:
        image_count = extract_images_from_pptx(filepath, output_dir)
    image_count = max(image_count, img_count_meta)

    # --- Step 4: Build manifest ---
    manifest = PptxManifest(
        source=os.path.basename(filepath),
        total_slides=len(slides),
        slides_with_notes=sum(1 for s in slides if s.has_notes),
        slides_with_tables=sum(1 for s in slides if s.has_table),
        total_images=image_count,
        slides=slides,
        issues=issues,
    )
    _write_manifest(manifest, output_dir)
    _print_summary(manifest)

    if verify:
        report = verify_pptx_quality(manifest)
        report_path = os.path.join(output_dir, "quality-report.json")
        with open(report_path, "w", encoding="utf-8") as f:
            json.dump(report, f, indent=2, ensure_ascii=False)
        print_pptx_quality(report)
        logger.info(f"Quality report -> {report_path}")

    return manifest


def _write_manifest(manifest: PptxManifest, output_dir: str) -> str:
    path = os.path.join(output_dir, "manifest.json")
    with open(path, "w", encoding="utf-8") as f:
        json.dump(asdict(manifest), f, indent=2, ensure_ascii=False)
    logger.info(f"Manifest → {path}")
    return path


def _print_summary(m: PptxManifest) -> None:
    print(f"""
═══════════════════════════════════════════
  PPTX Preprocessing Complete
═══════════════════════════════════════════
  Source:          {m.source}
  Slides:          {m.total_slides}
  With notes:      {m.slides_with_notes}
  With tables:     {m.slides_with_tables}
  Images:          {m.total_images}
  Issues:          {len(m.issues)}
═══════════════════════════════════════════
""")


def process_batch(input_dir: str, output_base: str) -> list[PptxManifest]:
    results: list[PptxManifest] = []
    files = list(Path(input_dir).glob("*.pptx"))
    logger.info(f"Batch processing {len(files)} PPTX file(s) in {input_dir}")
    for f in files:
        output_dir = os.path.join(output_base, f.stem)
        logger.info(f"--- Processing: {f.name} ---")
        results.append(process_pptx(str(f), output_dir))
    return results


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main() -> None:
    parser = argparse.ArgumentParser(
        description="PPTX Preprocessing for SE Requirements Decomposition",
    )
    parser.add_argument("--input", "-i", required=True)
    parser.add_argument("--output", "-o", required=True)
    parser.add_argument("--batch", action="store_true")
    parser.add_argument("--no-images", action="store_true")
    parser.add_argument("--verify", action="store_true", help="Run quality assessment on extraction results")
    parser.add_argument("--verbose", "-v", action="store_true")

    args = parser.parse_args()
    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)

    if args.batch:
        if not os.path.isdir(args.input):
            print(f"ERROR: --batch requires a directory", file=sys.stderr)
            sys.exit(1)
        process_batch(args.input, args.output)
    else:
        if not os.path.isfile(args.input):
            print(f"ERROR: File not found: {args.input}", file=sys.stderr)
            sys.exit(1)
        process_pptx(args.input, args.output, extract_images=not args.no_images, verify=args.verify)


if __name__ == "__main__":
    main()
